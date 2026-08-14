from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from processing.cleaner import TextCleaner

from loaders.base_loader import BaseLoader


class PPTXLoader(BaseLoader):

    def __init__(self, file_path):

        self.file_path = Path(file_path)

    def _extract_shape_text(self, shape, texts):

        # Grouped shapes don't expose .text directly - recurse into them.
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:

            for sub_shape in shape.shapes:
                self._extract_shape_text(sub_shape, texts)

            return

        # Tables were previously dropped entirely.
        if getattr(shape, "has_table", False):

            for row in shape.table.rows:

                row_text = " | ".join(
                    cell.text for cell in row.cells
                )

                if row_text.strip():
                    texts.append(row_text)

            return

        if hasattr(shape, "text"):
            texts.append(shape.text)

    def load(self):

        if not self.file_path.exists():

            raise FileNotFoundError(
                f"{self.file_path} not found."
            )

        presentation = Presentation(
            self.file_path
        )

        slides = []

        for number, slide in enumerate(
            presentation.slides,
            start=1
        ):

            texts = []

            for shape in slide.shapes:
                self._extract_shape_text(shape, texts)

            slides.append({

                "page_number": number,

                "text": TextCleaner.clean(
                    " ".join(texts)
                ),

                "metadata": {

                    "source": self.file_path.name,

                    "page": number,

                    "file_type": "pptx"

                }

            })

        return slides
